"""
generate_pptx.py
Builds a Ledger Strategy Review deck from the live dashboard data dict.
Called by the Flask /export/pptx route.
"""

from __future__ import annotations
import tempfile
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE


# ── Palette ──────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1E, 0x27, 0x61)
INDIGO  = RGBColor(0x4F, 0x46, 0xE5)
INDIGO_L= RGBColor(0x81, 0x8C, 0xF8)
TEAL    = RGBColor(0x0D, 0x94, 0x88)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE=RGBColor(0xF8, 0xFA, 0xFC)
SLATE   = RGBColor(0x64, 0x74, 0x8B)
DARK    = RGBColor(0x0F, 0x17, 0x2A)
GREEN   = RGBColor(0x10, 0xB9, 0x81)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
RED     = RGBColor(0xEF, 0x44, 0x44)
PURPLE  = RGBColor(0x8B, 0x5C, 0xF6)
LIGHT_GRAY = RGBColor(0xE2, 0xE8, 0xF0)

STATUS_COLORS = {
    "very-high":  GREEN,
    "validated":  RGBColor(0x22, 0xC5, 0x5E),
    "upgraded":   INDIGO,
    "medium":     AMBER,
    "downgraded": RGBColor(0xF9, 0x73, 0x16),
    "hypothesis": PURPLE,
    "rejected":   SLATE,
}

STAT_COLORS = [INDIGO, AMBER, RED, PURPLE]
OKR_COLORS  = [INDIGO, GREEN, AMBER, TEAL, PURPLE, RGBColor(0xEC,0x48,0x99), RGBColor(0xF9,0x73,0x16)]


# ── Low-level helpers ─────────────────────────────────────────────────────────

def _no_border(shape):
    shape.line.fill.background()


def _rect(slide, x, y, w, h, fill_rgb: RGBColor, border=False):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if not border:
        _no_border(shape)
    return shape


def _text(slide, text, x, y, w, h, *, size=12, bold=False, italic=False,
          color=DARK, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
          wrap=True, font="Calibri"):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap   = wrap
    tf.auto_size   = None
    tf.vertical_anchor = valign
    para = tf.paragraphs[0]
    para.alignment = align
    run  = para.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic= italic
    run.font.color.rgb = color
    return txb


def _header(slide, section_label, subtitle=""):
    """Dark navy header bar for content slides."""
    _rect(slide, 0, 0, 10, 0.85, NAVY)
    _text(slide, section_label, 0.45, 0.0, 5, 0.85,
          size=12, bold=True, color=WHITE, valign=MSO_ANCHOR.MIDDLE,
          font="Calibri")
    if subtitle:
        _text(slide, subtitle, 0, 0.0, 9.6, 0.85,
              size=11, color=INDIGO_L, align=PP_ALIGN.RIGHT,
              valign=MSO_ANCHOR.MIDDLE, font="Calibri")


# ── Slide builders ────────────────────────────────────────────────────────────

def _slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

    # Left accent bar
    _rect(slide, 0, 0, 0.22, 5.625, INDIGO)

    # Label
    _text(slide, "PRODUCT BRAIN", 0.5, 1.35, 9, 0.3,
          size=10, bold=True, color=INDIGO_L, font="Calibri")

    # Title
    _text(slide, "Ledger Strategy Review", 0.5, 1.7, 9, 1.1,
          size=42, bold=True, color=WHITE, font="Calibri")

    # Subtitle
    _text(slide, "20 customer interviews · 9 strategic bets · 4 competitors · 6 market shifts",
          0.5, 2.85, 9, 0.45, size=13, color=INDIGO_L, font="Calibri")

    # Date — anchored just below subtitle with more breathing room
    today = datetime.now().strftime("%B %Y")
    _text(slide, today, 0.5, 3.45, 9, 0.3, size=11, color=INDIGO_L, font="Calibri")


def _slide_market(prs, market_stats):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = OFF_WHITE

    _header(slide, "MARKET CONTEXT", "6 macro shifts shaping the opportunity")

    positions = [(0.4, 1.05), (5.15, 1.05), (0.4, 3.1), (5.15, 3.1)]
    for i, stat in enumerate(market_stats[:4]):
        x, y = positions[i]
        color = STAT_COLORS[i % len(STAT_COLORS)]

        # Card background
        _rect(slide, x, y, 4.45, 1.8, WHITE)
        # Accent left bar
        _rect(slide, x, y, 0.07, 1.8, color)

        # Value
        _text(slide, stat["value"], x + 0.2, y + 0.18, 4.0, 0.65,
              size=30, bold=True, color=color, font="Calibri")
        # Label
        _text(slide, stat["label"], x + 0.2, y + 0.82, 4.0, 0.32,
              size=12, bold=True, color=DARK, font="Calibri")
        # Sub
        _text(slide, stat["sub"], x + 0.2, y + 1.17, 4.0, 0.32,
              size=10, color=SLATE, font="Calibri")


def _slide_bets(prs, bets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = OFF_WHITE

    _header(slide, "STRATEGIC BETS", "Updated confidence · scale 0–5")

    # Horizontal bar chart
    chart_data = ChartData()
    chart_data.categories = [b["name"] for b in bets]
    chart_data.add_series("Confidence", [b["confidence"] for b in bets])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.3), Inches(0.95), Inches(9.4), Inches(4.55),
        chart_data
    )
    chart = chart_frame.chart

    # Style the chart
    chart.chart_style = 2
    plot = chart.plots[0]
    plot.gap_width = 60

    # Color each bar by status
    series = chart.series[0]
    for idx, bet in enumerate(bets):
        pt = series.points[idx]
        color = STATUS_COLORS.get(bet["status"], INDIGO)
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = color

    # Value axis max
    chart.value_axis.maximum_scale = 5.0
    chart.value_axis.minimum_scale = 0.0
    chart.has_legend = False

    # Axis label font
    from pptx.util import Pt as Pt2
    chart.value_axis.tick_labels.font.size  = Pt2(9)
    chart.category_axis.tick_labels.font.size = Pt2(9)


def _slide_competitors(prs, competitors):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = OFF_WHITE

    _header(slide, "COMPETITIVE LANDSCAPE", "Direct competitors · Tier 1")

    THREAT_COLOR = {"HIGH": RED, "MEDIUM": AMBER, "LOW": GREEN, "LOW-MEDIUM": TEAL}

    # Column headers
    headers = ["Competitor", "Funding", "Customers", "Threat Level"]
    col_x   = [0.4, 3.6, 5.55, 7.5]
    col_w   = [3.0, 1.8, 1.8, 2.0]

    # Header row background
    _rect(slide, 0.4, 1.0, 9.2, 0.45, NAVY)

    for j, h in enumerate(headers):
        _text(slide, h, col_x[j] + 0.1, 1.0, col_w[j], 0.45,
              size=10, bold=True, color=WHITE, valign=MSO_ANCHOR.MIDDLE, font="Calibri")

    # Data rows
    for i, comp in enumerate(competitors):
        row_y  = 1.48 + i * 0.72
        row_bg = WHITE if i % 2 == 0 else RGBColor(0xF1, 0xF5, 0xF9)
        _rect(slide, 0.4, row_y, 9.2, 0.68, row_bg)

        threat_col = THREAT_COLOR.get(comp["threat"], AMBER)
        funding_str = f"${comp['funding']}M raised" if comp["funding"] > 0 else "Bootstrapped"
        customers_str = f"~{comp['customers']:,}"

        row_vals = [comp["name"], funding_str, customers_str, comp["threat"]]
        row_colors = [DARK, SLATE, SLATE, threat_col]
        row_bold   = [True, False, False, True]

        for j in range(4):
            _text(slide, row_vals[j], col_x[j] + 0.1, row_y, col_w[j], 0.68,
                  size=10, bold=row_bold[j], color=row_colors[j],
                  valign=MSO_ANCHOR.MIDDLE, font="Calibri")


def _slide_okrs(prs, okr_targets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = OFF_WHITE

    _header(slide, "2024 OKR TARGETS", "Pre-launch baseline → year-end targets")

    for i, okr in enumerate(okr_targets[:6]):
        col = i % 2
        row = i // 2
        x = 0.4  + col * 4.85
        y = 1.05 + row * 1.45

        color_hex = okr.get("color", "#6366f1").lstrip("#")
        try:
            r, g, b = int(color_hex[0:2],16), int(color_hex[2:4],16), int(color_hex[4:6],16)
            color = RGBColor(r, g, b)
        except Exception:
            color = INDIGO

        _rect(slide, x, y, 4.5, 1.2, WHITE)
        _rect(slide, x, y, 0.07, 1.2, color)

        _text(slide, okr["label"], x + 0.22, y + 0.1, 3.5, 0.35,
              size=12, bold=True, color=DARK, font="Calibri")

        unit = okr["unit"] if okr["unit"] else ""
        target_val = int(okr["target"]) if okr["target"] == int(okr["target"]) else okr["target"]
        target_str = f"Target: {target_val}{unit}"
        _text(slide, target_str, x + 0.22, y + 0.5, 3.0, 0.3,
              size=11, color=SLATE, font="Calibri")

        _text(slide, "PRE-LAUNCH", x + 3.1, y + 0.1, 1.3, 0.3,
              size=9, bold=True, color=SLATE, align=PP_ALIGN.RIGHT, font="Calibri")


def _slide_recommendations(prs, bets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

    # Left accent bar
    _rect(slide, 0, 0, 0.22, 5.625, TEAL)

    _text(slide, "RECOMMENDATIONS", 0.5, 0.4, 9, 0.3,
          size=10, bold=True, color=TEAL, font="Calibri")
    _text(slide, "Top bets to prioritise now", 0.5, 0.72, 9, 0.7,
          size=30, bold=True, color=WHITE, font="Calibri")

    # Top 4 bets by confidence, non-rejected
    top = sorted(
        [b for b in bets if b["status"] not in ("rejected", "hypothesis")],
        key=lambda b: b["confidence"], reverse=True
    )[:4]

    for i, bet in enumerate(top):
        y = 1.65 + i * 0.97

        # Card (semi-transparent white — simulate with light navy)
        _rect(slide, 0.5, y, 9.0, 0.82, RGBColor(0x2A, 0x36, 0x7A))
        _rect(slide, 0.5, y, 0.07, 0.82, TEAL)

        # Rank number
        _text(slide, f"#{i+1}", 0.7, y + 0.05, 0.45, 0.35,
              size=16, bold=True, color=TEAL, font="Calibri")

        # Bet name
        _text(slide, bet["name"], 1.25, y + 0.05, 7.0, 0.35,
              size=13, bold=True, color=WHITE, font="Calibri")

        # Evidence / confidence line
        evidence = bet.get("evidence", 0)
        evidence_str = f"{evidence} interviews" if evidence else "see hypothesis-evidence.md"
        meta = f"Confidence: {bet['label']}  ·  Evidence: {evidence_str}"
        _text(slide, meta, 1.25, y + 0.45, 7.0, 0.27,
              size=10, color=WHITE, font="Calibri")


# ── Public entry point ────────────────────────────────────────────────────────

def build_pptx(data: dict) -> Path:
    """Generate the deck and return a Path to the temp .pptx file."""
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    _slide_title(prs)
    _slide_market(prs, data.get("market_stats", []))
    _slide_bets(prs, data.get("bets", []))
    _slide_competitors(prs, data.get("competitors", []))
    _slide_okrs(prs, data.get("okr_targets", []))
    _slide_recommendations(prs, data.get("bets", []))

    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    tmp.close()
    prs.save(tmp.name)
    return Path(tmp.name)
